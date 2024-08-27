from .pii import PII
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from io import StringIO
from .const import ASSETS

class PPTX(PII):
    def __init__(self, file_loc: str):
        self.ppt = Presentation(file_loc)
        self.pages = len(self.ppt.slides)
    
    def save_image(self, shape):
        image = shape.image
        image_bytes = image.blob
        with open(f"{ASSETS}/{shape.shape_id}.png", "wb") as fp:
            fp.write(image_bytes)

    def get_metadata(self):
        props = self.ppt.core_properties
        result = {
            "info": 
            {
                "title": props.title,
                "author": props.author,
                "last_modified_by": props.last_modified_by,
                "modified": props.modified,
                "created": props.created
            },
            "total_pages": self.pages
        }
        return result

    def read_all(self, save_images: bool=True) -> str:
        result = []
        for idx, slide in enumerate(self.ppt.slides):
            tmp = {"page": idx, "text": [], "images": []}
            for shape in slide.shapes:
                if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                    if save_images == True:
                        tmp["images"].append(shape.shape_id)
                        self.save_image(shape)
                if not shape.has_text_frame:
                    continue
                for text_frame in shape.text_frame.paragraphs:
                    tmp["text"].append(text_frame.text)
            result.append(tmp)
        return result
    
    def read_page(self, page: int, save_images: bool=True) -> str:
        if page not in range(self.pages):
            return None
        
        result = {"page": page, "text": [], "images": []}
        for shape in self.ppt.slides[page].shapes:
            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                if save_images == True:
                    result["images"].append(shape.shape_id)
                    self.save_image(shape)
            if not shape.has_text_frame:
                continue
            for text_frame in shape.text_frame.paragraphs:
                result["text"].append(text_frame.text)
        return result

                    